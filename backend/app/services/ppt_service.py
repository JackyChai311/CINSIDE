"""PPT Section 拆并核心服务。

移植自 Hansonus 项目的 pptsections-svc.ts，适配 CINSIDE 的 Python 后端。

核心流程：
1. analyze_files — 解析每个 PPT 的幻灯片及文本节点
2. detect_sections — AI 识别 section 章节
3. build_reading_script — 生成 Markdown 阅读稿
4. build_merged_ppt — 合并各 section 为一个总览 PPT
5. modify_content — AI 统一修改内容，输出补丁
6. apply_patches — 一键回填到原始 PPT
"""
from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
import base64
import asyncio
from pathlib import Path
from typing import Any, Callable

import httpx

from ..config import settings
from . import image_gen
from . import officecli

# MarkItDown 快速文本提取（可选依赖）
try:
    from markitdown import MarkItDown
    _MD = MarkItDown()
    _HAS_MARKITDOWN = True
except Exception:
    _MD = None
    _HAS_MARKITDOWN = False

# ---- 类型定义（dataclass 风格的 dict） ----
# TextNode = {"path": str, "text": str}
# SlideNode = {"index": int, "path": str, "title": str, "texts": list[TextNode]}
# FileSlides = {"file_id": str, "file_name": str, "file_path": str, "slides": list[SlideNode]}
# SectionPart = {"file_id": str, "file_name": str, "file_path": str,
#                "slide_start": int, "slide_end": int, "slides": list[SlideNode]}
# Section = {"name": str, "parts": list[SectionPart]}
# TextPatch = {"file_id": str, "file_name": str, "file_path": str,
#              "slide": int, "path": str, "text": str, "new_text": str}


# ---- 临时目录 ----

def _get_work_dir() -> Path:
    """获取 PPT 工作目录（存放合并文件、截图等）。"""
    work = Path(tempfile.gettempdir()) / "cinside-ppt"
    work.mkdir(parents=True, exist_ok=True)
    return work


def _make_file_id(file_path: str) -> str:
    """从文件路径生成简短的 file_id。"""
    return str(abs(hash(file_path)))


def get_page_count(file_path: str) -> int:
    """获取参考文件页数（PPT 用 OfficeCLI outline，PDF 用 PyMuPDF）。"""
    file_path = str(Path(file_path).resolve())
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    ext = Path(file_path).suffix.lower()
    if ext in (".ppt", ".pptx"):
        try:
            outline = officecli.view_outline(file_path)
            slides = _parse_outline_slides(outline, file_path)
            if slides:
                return len(slides)
        except Exception:
            pass
        return 1
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF
            with fitz.open(file_path) as doc:
                return max(doc.page_count, 1)
        except Exception:
            return 1
    return 1


# ---- 递归收集文本节点 ----

def _collect_text_nodes(node: Any, out: list[dict] | None = None) -> list[dict]:
    """递归遍历 JSON 节点树，收集 paragraph/table-cell 文本节点。"""
    if out is None:
        out = []
    if not node or not isinstance(node, dict):
        return out

    node_type = node.get("type", "")
    is_text_leaf = node_type in ("paragraph", "table-cell")
    if is_text_leaf:
        text = node.get("text", "")
        if isinstance(text, str) and text.strip():
            out.append({"path": node.get("path", ""), "text": text})

    children = node.get("children", [])
    if isinstance(children, list):
        for child in children:
            _collect_text_nodes(child, out)
    return out


# ---- 步骤 1：解析每个 PPT 的幻灯片及文本 ----

# Progress callback: (current_index, total, file_name, status, extra_dict)
ProgressCb = Callable[[int, int, str, str, dict], None]


def analyze_files(
    file_paths: list[str],
    progress_cb: ProgressCb | None = None,
) -> list[dict]:
    """解析多个 PPT 文件，返回每个文件的幻灯片及文本节点列表。

    优化流程：
    1. 优先用 MarkItDown 一次转换整份 PPT 为 markdown，从中同时提取
       幻灯片数量、标题和正文（无需额外调用 OfficeCLI view_outline）。
    2. MarkItDown 失败或未给出分页标记时，回退到 OfficeCLI view_outline。
    3. 个别页 MarkItDown 未覆盖到的，再用 OfficeCLI get_node 逐页补取。
    """
    def _report(idx: int, total: int, name: str, status: str, **extra: Any) -> None:
        if progress_cb:
            try:
                progress_cb(idx, total, name, status, extra)
            except Exception:
                pass

    result: list[dict] = []
    total = len(file_paths)
    print(f"[PPT-ANALYZE] 开始解析 {total} 个 PPT 文件 (markitdown={_HAS_MARKITDOWN})", flush=True)

    for i, fp in enumerate(file_paths):
        file_path = str(Path(fp).resolve())
        file_name = Path(file_path).name

        if not os.path.exists(file_path):
            _report(i + 1, total, file_name, "error", message="文件不存在")
            raise FileNotFoundError(f"文件不存在: {file_path}")

        file_id = _make_file_id(file_path)
        _report(i, total, file_name, "parsing")
        print(f"[PPT-ANALYZE] [{i+1}/{total}] {file_name}", flush=True)

        # 1) MarkItDown 快速提取全文 markdown
        md_text = ""
        if _HAS_MARKITDOWN:
            try:
                _report(i, total, file_name, "markitdown")
                md_result = _MD.convert(file_path)
                md_text = md_result.text_content or ""
                print(f"[PPT-ANALYZE]   MarkItDown: {len(md_text)} chars", flush=True)
            except Exception as e:
                print(f"[PPT-ANALYZE]   MarkItDown 失败，将用 OfficeCLI: {e}", flush=True)
                md_text = ""

        # 2) 从 markdown 中解析幻灯片大纲（数量 + 标题）和文本
        slides_meta: list[dict] = []
        slide_texts_map: dict[int, list[str]] = {}

        if md_text:
            slides_meta, slide_texts_map = _parse_markdown_slides(md_text)

        # 3) MarkItDown 未能提供幻灯片结构时，回退 OfficeCLI view_outline
        if not slides_meta:
            _report(i, total, file_name, "outline")
            outline = officecli.view_outline(file_path)
            slides_meta = _parse_outline_slides(outline, file_path)
            print(f"[PPT-ANALYZE]   OfficeCLI outline: {len(slides_meta)} slides", flush=True)
            # markdown 有文本但没分页标记时，把全部文本归到第 1 页
            if md_text and not slide_texts_map:
                lines = [ln.strip() for ln in md_text.splitlines() if ln.strip()]
                if lines:
                    slide_texts_map[1] = lines

        _report(i, total, file_name, "outline", slides=len(slides_meta))

        # 4) 组装 slides；MarkItDown 未覆盖到的页用 OfficeCLI 逐页补取
        slides: list[dict] = []
        for s in slides_meta:
            idx = s["index"]
            texts: list[dict] = []

            # 优先用 MarkItDown 的文本
            if idx in slide_texts_map and slide_texts_map[idx]:
                for line in slide_texts_map[idx]:
                    if line.strip():
                        texts.append({"path": f"/slide[{idx}]/md", "text": line.strip()})
            else:
                # 回退：OfficeCLI 逐页 get_node
                try:
                    node_json = officecli.get_node(file_path, f"/slide[{idx}]")
                    slide_node = _extract_slide_node(node_json)
                    if slide_node:
                        texts = _collect_text_nodes(slide_node)
                except Exception as e:
                    print(f"[PPT-ANALYZE]   警告: slide[{idx}] 解析失败: {e}", flush=True)
                    texts = []

            slides.append({
                "index": idx,
                "path": f"/slide[{idx}]",
                "title": s.get("title", ""),
                "texts": texts,
            })

        total_texts = sum(len(s["texts"]) for s in slides)
        print(f"[PPT-ANALYZE]   -> {len(slides)} 张幻灯片, {total_texts} 个文本节点", flush=True)

        result.append({
            "file_id": file_id,
            "file_name": file_name,
            "file_path": file_path,
            "slides": slides,
        })
        _report(i + 1, total, file_name, "done", slides=len(slides), texts=total_texts)

    print("[PPT-ANALYZE] 全部文件解析完成", flush=True)
    return result


def _parse_markdown_slides(md_text: str) -> tuple[list[dict], dict[int, list[str]]]:
    """从 MarkItDown 输出的 markdown 中同时解析幻灯片大纲和文本。

    返回 (slides_meta, slide_texts_map)：
    - slides_meta: [{"index": int, "title": str}, ...]
    - slide_texts_map: {slide_num: [line, ...]}

    MarkItDown 的 pptx 转换器通常用 ``<!-- Slide number: N -->`` 注释分页，
    也可能用 ``## Slide N`` 之类的标题。我们兼容多种标记。
    """
    if not md_text:
        return [], {}

    # 匹配幻灯片分隔标记
    # 形式1: <!-- Slide number: N -->
    # 形式2: <!-- Slide N -->
    # 形式3: ## Slide N  / # Slide N
    slide_pattern = re.compile(
        r"(?:<!--\s*Slide\s*(?:number\s*:\s*)?(\d+)\s*-->|"
        r"^#{1,3}\s*Slide\s+(\d+))",
        re.IGNORECASE | re.MULTILINE,
    )

    matches = list(slide_pattern.finditer(md_text))
    if not matches:
        # 没有明确分页标记，无法确定幻灯片数量，返回空让调用方回退
        return [], {}

    slides_meta: list[dict] = []
    slide_texts_map: dict[int, list[str]] = {}

    for mi, match in enumerate(matches):
        slide_num = int(match.group(1) or match.group(2))
        start = match.end()
        end = matches[mi + 1].start() if mi + 1 < len(matches) else len(md_text)
        chunk = md_text[start:end]

        lines: list[str] = []
        title = ""
        for ln in chunk.splitlines():
            ln = ln.strip()
            if not ln:
                continue
            # 跳过纯 markdown 符号行
            if re.match(r"^[-*#=]+$", ln):
                continue
            # 第一个看起来像标题的行作为幻灯片标题
            if not title:
                heading = re.match(r"^#{1,6}\s+(.+)$", ln)
                if heading:
                    title = heading.group(1).strip()
                elif len(lines) == 0 and len(ln) <= 80:
                    # 第一行且不太长，暂作为标题候选
                    title = ln
            lines.append(ln)

        # 标题也加入文本列表（保持完整内容）
        if lines:
            slide_texts_map[slide_num] = lines
        slides_meta.append({"index": slide_num, "title": title})

    return slides_meta, slide_texts_map


def _split_markdown_by_slide(md_text: str, expected_count: int) -> dict[int, list[str]]:
    """把 MarkItDown 输出的 markdown 按幻灯片切分成文本行列表（兼容保留）。"""
    _, texts = _parse_markdown_slides(md_text)
    return texts


def _parse_outline_slides(outline_json: dict, file_path: str) -> list[dict]:
    """从 outline JSON 或文本中解析幻灯片列表。

    优先从 JSON 结构提取；失败则回退到文本正则。
    """
    # 尝试从 JSON 结构提取
    if isinstance(outline_json, dict):
        # 可能的结构: {"slides": [...]} 或 {"data": {"slides": [...]}}
        slides_data = outline_json.get("slides") or outline_json.get("data", {}).get("slides")
        if isinstance(slides_data, list):
            result = []
            for s in slides_data:
                if isinstance(s, dict):
                    result.append({
                        "index": s.get("index", s.get("number", 0)),
                        "title": s.get("title", s.get("name", "")),
                    })
            if result:
                return result

    # 回退：用文本格式 + 正则解析
    try:
        result = officecli._run(["view", file_path, "outline"], timeout=30)
        text = result.stdout if result.returncode == 0 else ""
    except Exception:
        text = ""

    slides: list[dict] = []
    for m in re.finditer(r"Slide\s+(\d+):\s*(.+)", text):
        slides.append({"index": int(m.group(1)), "title": m.group(2).strip()})

    if not slides:
        # 如果大纲解析失败，尝试获取幻灯片数量
        try:
            stats = officecli._run_json(["view", file_path, "stats", "--json"], timeout=15)
            count = stats.get("slides") or stats.get("slideCount") or stats.get("data", {}).get("slides")
            if isinstance(count, int) and count > 0:
                slides = [{"index": i + 1, "title": ""} for i in range(count)]
        except Exception:
            pass

    return slides


def _extract_slide_node(node_json: dict) -> dict | None:
    """从 get --json 的输出中提取幻灯片节点。"""
    if not isinstance(node_json, dict):
        return None
    # 结构: {"data": {"results": [node]}}
    results = node_json.get("data", {}).get("results")
    if isinstance(results, list) and results:
        return results[0] if isinstance(results[0], dict) else None
    # 也可能直接是节点本身
    if node_json.get("type") == "slide":
        return node_json
    return None


# ---- 步骤 2：AI 识别 section ----

async def detect_sections(file_slides_list: list[dict], instruction: str = "") -> list[dict]:
    """调用 LLM 分析每个 PPT 的幻灯片，归类到 section。

    instruction 为用户的自然语言指令，例如"帮我把读课文部分拆出来"，
    AI 会据此调整识别重点和归类策略。
    """
    sections: list[dict] = []
    print(f"[PPT-DETECT] 开始 AI 识别 section: {len(file_slides_list)} 个文件, instruction={instruction!r}", flush=True)

    for i, fs in enumerate(file_slides_list):
        system_prompt = (
            "你是 PPT 结构分析师。用户提供一批\"同模板\"PPT，每个 PPT 内部由若干逻辑 section（章节）组成，"
            "同一 section 在不同 PPT 中的页数可能不同。\n"
            "请把该 PPT 的每一页归类到某个 section，section 名称用简洁的中文（各 PPT 相同含义的 section 应使用完全一致的名称）。\n"
            "根据幻灯片标题和内容判断归属；相邻页通常属于同一 section，尽量给出连续的范围。\n"
        )
        if instruction.strip():
            system_prompt += (
                f"\n用户特别要求：{instruction.strip()}\n"
                "请优先按照用户的要求来识别和归类章节；如果用户要求提取/拆分某些特定部分，"
                "请确保将相关幻灯片单独归为对应的 section，其余部分也应合理归类。\n"
            )
        system_prompt += (
            "严格以 JSON 返回，格式：{\"sections\":[{\"name\":\"客户背景\",\"slides\":[1,2,3]}]}，"
            "其中 slides 是该 section 的幻灯片序号数组。不要输出任何其他文字。"
        )

        user_content = f"文件名：{fs['file_name']}\n\n幻灯片列表：\n{_slide_summary(fs)}"

        try:
            raw_sections = await _call_llm_json(system_prompt, user_content, temperature=0.2)
            print(f"[PPT-DETECT] [{i+1}/{len(file_slides_list)}] {fs['file_name']} -> "
                  f"AI 识别出 {len(raw_sections)} 个 section", flush=True)

            for raw in raw_sections:
                name = str(raw.get("name", "")).strip()
                if not name:
                    continue
                slide_nums = raw.get("slides", [])
                if not isinstance(slide_nums, list):
                    continue
                slides_sorted = sorted(set(int(n) for n in slide_nums if isinstance(n, (int, float))))
                if not slides_sorted:
                    continue

                part_slides = [s for s in fs["slides"] if s["index"] in slides_sorted]
                part = {
                    "file_id": fs["file_id"],
                    "file_name": fs["file_name"],
                    "file_path": fs["file_path"],
                    "slide_start": slides_sorted[0],
                    "slide_end": slides_sorted[-1],
                    "slides": part_slides,
                }
                _push_section_part(sections, name, part)
                print(f"[PPT-DETECT]   . \"{name}\" -> slides {slides_sorted[0]}-{slides_sorted[-1]} "
                      f"({len(slides_sorted)} 页)", flush=True)
        except Exception as e:
            print(f"[PPT-DETECT] 错误: {fs['file_name']} 识别失败: {e}", flush=True)

    sec_summary = ", ".join(f"{s['name']}({len(s['parts'])}个片段)" for s in sections)
    print(f"[PPT-DETECT] 识别完成，共 {len(sections)} 个全局 section: {sec_summary}", flush=True)
    return sections


def _slide_summary(file_slides: dict) -> str:
    """生成幻灯片摘要文本供 LLM 分析。"""
    lines: list[str] = []
    for s in file_slides["slides"]:
        sample = " / ".join(t["text"] for t in s["texts"])
        sample = re.sub(r"\s+", " ", sample)[:120]
        title = s.get("title") or "(无标题)"
        lines.append(f"Slide {s['index']}: {title}" + (f" | {sample}" if sample else ""))
    return "\n".join(lines)


def _push_section_part(sections: list[dict], name: str, part: dict) -> None:
    """将 section part 添加到对应名称的 section 中（不存在则新建）。"""
    for s in sections:
        if s["name"] == name:
            s["parts"].append(part)
            return
    sections.append({"name": name, "parts": [part]})


# ---- 步骤 3：合并总览 ----

def build_reading_script(sections: list[dict]) -> str:
    """生成 Markdown 格式的阅读稿。"""
    lines = ["# PPT Sections 合并总览", ""]
    total_parts = sum(len(s["parts"]) for s in sections)
    lines.append(f"共 {len(sections)} 个 section，{total_parts} 个文件片段。")
    lines.append("")

    for section in sections:
        lines.append(f"## ▍{section['name']}")
        lines.append("")
        for part in section["parts"]:
            lines.append(f"### {part['file_name']}")
            lines.append(f"> 页面范围：第 {part['slide_start']}–{part['slide_end']} 页")
            lines.append("")
            for slide in part["slides"]:
                title = slide.get("title") or "(无标题)"
                lines.append(f"**第 {slide['index']} 页 · {title}**")
                for t in slide["texts"]:
                    clean = re.sub(r"\s+", " ", t["text"]).strip()
                    if clean:
                        lines.append(f"- {clean}")
                lines.append("")
    return "\n".join(lines)


# ---- 步骤 0：按文字新建 PPT ----

async def create_from_text(text: str) -> dict:
    """根据用户输入的文字（大纲/要点/正文）生成一份新 PPT。

    流程：
    1. 调用 LLM 把文字拆分成多页幻灯片大纲（每页含标题 + 要点）。
    2. 用 officecli 创建空白 PPT，逐页添加 slide 标题 + 正文 shape。
    3. 保存并返回生成文件信息。
    """
    print(f"[PPT-CREATE] 开始按文字新建 PPT: 输入 {len(text)} 字符", flush=True)

    system_prompt = (
        "你是专业 PPT 大纲设计师。用户给出一段文字或一个主题，请把它组织成一份结构清晰、可直接用于演示的 PPT 大纲。\n"
        "要求：\n"
        "1. 每页只聚焦一个主题，页数 4~12 页，内容较多时适当拆分；若用户给的文字本身较长，应充分展开，而不是简略带过。\n"
        "2. 每页包含一个有力的标题（title）和 4~6 条要点（bullets）。\n"
        "3. 每条要点必须是完整、自成一体的陈述句，能独立说明一个关键信息，尽量具体、有实质内容；\n"
        "   不要用含糊的短语或关键词堆砌，要让观众看着要点就能理解这一页在讲什么。\n"
        "4. 尊重用户原文的信息和语气，不要编造原文没有的事实；可以合理补充过渡性、解释性的说明。\n"
        "5. 中文输出。\n"
        "严格以 JSON 返回，格式：{\"slides\":[{\"title\":\"目录\",\"bullets\":[\"点1\",\"点2\"]}]}。"
        "不要输出任何其他文字。"
    )
    user_content = text

    try:
        raw = await _call_llm_json(system_prompt, user_content, temperature=0.4)
    except Exception as e:
        raise RuntimeError(f"AI 生成 PPT 大纲失败: {e}")

    slides = []
    for item in raw:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title", "")).strip()
        bullets = item.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = []
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        # 标题为空时用第一条要点兜底
        if not title and bullets:
            title = bullets.pop(0)
        if not title:
            continue
        slides.append({"title": title, "bullets": bullets})

    if not slides:
        raise RuntimeError("AI 未能生成有效的 PPT 大纲")

    work_dir = _get_work_dir()
    out_path = work_dir / f"新建-{int(time.time()*1000)}.pptx"
    print(f"[PPT-CREATE] 生成 {len(slides)} 页大纲 -> {out_path}", flush=True)

    if not officecli.create_document(str(out_path)):
        raise RuntimeError("创建新 PPT 失败")

    for i, slide in enumerate(slides):
        slide_num = i + 1
        # 添加标题
        officecli.add_element(str(out_path), "/", "slide", {"title": slide["title"]})
        # 添加正文（多条要点用换行拼成多段）
        if slide["bullets"]:
            body = "\n".join(slide["bullets"])
            officecli.add_element(
                str(out_path), f"/slide[{slide_num}]", "shape", {"text": body}
            )
        print(f"[PPT-CREATE]   slide[{slide_num}] \"{slide['title']}\" "
              f"({len(slide['bullets'])} 条要点)", flush=True)

    officecli.save_document(str(out_path))

    final_name = out_path.name
    print(f"[PPT-CREATE] 完成: {final_name} ({len(slides)} 页)", flush=True)
    return {
        "file_path": str(out_path),
        "file_name": final_name,
        "total_slides": len(slides),
    }


# ---- 版式渲染辅助（风格驱动） ----

class _ShapeCounter:
    """按添加顺序分配 shape[N] 路径，装饰开关不同时索引也不乱。"""

    def __init__(self, start: int = 2) -> None:
        self.i = start

    def add(self, out: str, sid: str, props: dict) -> str:
        officecli.add_element(out, sid, "shape", props)
        path = f"{sid}/shape[{self.i}]"
        self.i += 1
        return path


def _place_cover(out: str, sid: str, slide: dict, st, accent: str, tint: str,
                 has_image: bool, decor_cb) -> None:
    """封面页版式。shape[1] 为标题占位符（add slide 时已创建），本函数从 shape[2] 开始。"""
    sc = _ShapeCounter(2)
    bullets = slide["bullets"]
    # 全幅浅底
    sc.add(out, sid, {"x": "0in", "y": "0in", "width": "13.33in", "height": "7.5in",
                      "fill": tint, "geometry": "rect"})
    decor_cb("背景底")
    if st.cover_layout == "split":
        # 左侧大色块 + 右侧配图/留白
        sc.add(out, sid, {"x": "0in", "y": "0in", "width": "5.2in", "height": "7.5in",
                          "fill": accent, "geometry": "rect"})
        decor_cb("侧边色块")
        officecli.set_element(out, f"{sid}/shape[1]", {
            "x": "0.7in", "y": "2.35in", "width": "3.9in", "height": "2.0in",
            "size": str(st.title_size + 4), "bold": "true", "color": "#FFFFFF",
            "align": "left", "valign": "middle",
        })
        sc.add(out, sid, {"x": "0.7in", "y": "4.55in", "width": "1.4in", "height": "0.09in",
                          "fill": "#FFFFFF", "geometry": "rect"})
        decor_cb("强调条")
        if bullets:
            p = sc.add(out, sid, {"text": "  ·  ".join(bullets[:4])})
            officecli.set_element(out, p, {
                "x": "0.7in", "y": "4.8in", "width": "3.9in", "height": "1.5in",
                "size": "14", "color": "#FFFFFF", "align": "left", "valign": "top", "lineSpacing": "1.4",
            })
        p = sc.add(out, sid, {"text": "CINSIDE · AI 生成"})
        officecli.set_element(out, p, {
            "x": "0.7in", "y": "6.95in", "width": "3.9in", "height": "0.3in",
            "size": "11", "color": "#FFFFFF", "align": "left",
        })
    else:
        # center：顶部色带 + 居中大标题（有图时上半部留给配图）
        if "top_bar" in st.decor:
            sc.add(out, sid, {"x": "0in", "y": "0in", "width": "13.33in", "height": "0.22in",
                              "fill": accent, "geometry": "rect"})
            decor_cb("顶部色带")
        if has_image:
            officecli.set_element(out, f"{sid}/shape[1]", {
                "x": "1.2in", "y": "4.35in", "width": "10.93in", "height": "1.2in",
                "size": str(st.title_size + 6), "bold": "true", "color": st.title_color,
                "align": "center", "valign": "middle",
            })
            sc.add(out, sid, {"x": "5.67in", "y": "5.7in", "width": "2in", "height": "0.09in",
                              "fill": accent, "geometry": "rect"})
            decor_cb("强调条")
            if bullets:
                p = sc.add(out, sid, {"text": "  ·  ".join(bullets[:4])})
                officecli.set_element(out, p, {
                    "x": "1.2in", "y": "5.95in", "width": "10.93in", "height": "0.85in",
                    "size": "15", "color": st.body_color,
                    "align": "center", "valign": "top", "lineSpacing": "1.35",
                })
        else:
            officecli.set_element(out, f"{sid}/shape[1]", {
                "x": "1.2in", "y": "2.3in", "width": "10.93in", "height": "1.7in",
                "size": str(st.title_size + 12), "bold": "true", "color": st.title_color,
                "align": "center", "valign": "middle",
            })
            sc.add(out, sid, {"x": "5.67in", "y": "4.15in", "width": "2in", "height": "0.09in",
                              "fill": accent, "geometry": "rect"})
            decor_cb("强调条")
            if bullets:
                p = sc.add(out, sid, {"text": "  ·  ".join(bullets[:6])})
                officecli.set_element(out, p, {
                    "x": "1.2in", "y": "4.5in", "width": "10.93in", "height": "1.4in",
                    "size": str(st.body_size), "color": st.body_color,
                    "align": "center", "valign": "top", "lineSpacing": "1.4",
                })
        p = sc.add(out, sid, {"text": "CINSIDE · AI 生成"})
        officecli.set_element(out, p, {
            "x": "1.2in", "y": "6.95in", "width": "10.93in", "height": "0.3in",
            "size": "11", "color": "#94A3B8", "align": "center",
        })


def _place_bullet_cards(
    sc: "_ShapeCounter", out: str, sid: str,
    bullets: list[str],
    x: str, y_start: float, width: str,
    accent: str, tint: str, body_color: str, body_size: int,
    card_gap: float = 0.14,
    card_fill: str = "#FFFFFF",
) -> list[str]:
    """将每条要点渲染为独立的圆角矩形卡片，返回文字 shape 路径列表。

    每张卡片包含：左侧主题色竖条 + 序号圆点 + 文字区域。
    卡片之间有间距，避免纯文字堆砌的视觉效果。
    """
    paths: list[str] = []
    n = len(bullets)
    if n == 0:
        return paths

    # 估算每张卡片高度（按字数估算行数）
    # PPT 16:9 画幅，正文宽度内每行约容纳 36~40 个中文字符
    chars_per_line = 34
    card_heights: list[float] = []
    for bullet in bullets:
        est_lines = max(1, (len(bullet) + chars_per_line - 1) // chars_per_line)
        # 单行卡片更紧凑，多行卡片按行高扩展
        card_heights.append(max(0.58, 0.38 + est_lines * 0.26))

    # 若总高度超出可用空间，按比例压缩
    total_h = sum(card_heights) + card_gap * (n - 1)
    available_h = 4.8
    if total_h > available_h:
        scale = available_h / total_h
        card_heights = [max(0.48, h * scale) for h in card_heights]

    # 逐张放置（累积 y 坐标，避免重叠）
    cur_y = y_start
    for i, bullet in enumerate(bullets):
        card_h = card_heights[i]
        y = f"{cur_y:.2f}in"
        h = f"{card_h:.2f}in"

        # 卡片底色（浅色圆角矩形，有背景图时半透明）
        sc.add(out, sid, {
            "x": x, "y": y, "width": width, "height": h,
            "fill": card_fill, "geometry": "roundRect",
            "adj": "adj:val 10000",
        })
        # 左侧主题色竖条（更醒目）
        sc.add(out, sid, {
            "x": x, "y": y, "width": "0.09in", "height": h,
            "fill": accent, "geometry": "roundRect",
            "adj": "adj:val 6000",
        })
        # 序号圆点（更大更清晰）
        num_d = "0.34in"
        num_x = f"{float(x.rstrip('in')) + 0.22:.2f}in"
        num_y = f"{float(y.rstrip('in')) + (card_h - 0.34) / 2:.2f}in"
        sc.add(out, sid, {
            "x": num_x, "y": num_y, "width": num_d, "height": num_d,
            "fill": accent, "geometry": "ellipse",
        })
        # 序号文字
        num_path = sc.add(out, sid, {"text": str(i + 1)})
        officecli.set_element(out, num_path, {
            "x": num_x, "y": num_y, "width": num_d, "height": num_d,
            "size": "12", "bold": "true", "color": "#FFFFFF",
            "align": "center", "valign": "middle",
        })
        # 要点文字（先留空，由流式循环逐张填充）
        text_x = f"{float(x.rstrip('in')) + 0.72:.2f}in"
        text_w = f"{float(width.rstrip('in')) - 0.95:.2f}in"
        text_path = sc.add(out, sid, {"text": ""})
        officecli.set_element(out, text_path, {
            "x": text_x, "y": y, "width": text_w, "height": h,
            "size": str(body_size), "color": body_color,
            "valign": "middle", "lineSpacing": "1.4",
            "margin": "0.04in",
        })
        paths.append(text_path)
        cur_y += card_h + card_gap

    return paths


def _place_content(out: str, sid: str, slide_num: int, slide: dict, st,
                   accent: str, tint: str, has_image: bool, decor_cb,
                   add_background: bool = False) -> list[str]:
    """内容页版式，返回每张要点卡片的文字 shape 路径列表（无正文返回空列表）。shape[1] 为标题占位符。"""
    sc = _ShapeCounter(2)
    bullets = slide["bullets"]
    bullet_paths: list[str] = []

    if st.content_layout == "sidebar":
        # 左侧大色块 + 白字标题 + 右侧正文卡片
        sc.add(out, sid, {"x": "0in", "y": "0in", "width": "4.2in", "height": "7.5in",
                          "fill": accent, "geometry": "rect"})
        decor_cb("侧边色块")
        officecli.set_element(out, f"{sid}/shape[1]", {
            "x": "0.55in", "y": "0.9in", "width": "3.2in", "height": "2.6in",
            "size": str(st.title_size - 2), "bold": "true", "color": "#FFFFFF", "valign": "top",
        })
        sc.add(out, sid, {"x": "0.55in", "y": "3.7in", "width": "1.2in", "height": "0.09in",
                          "fill": "#FFFFFF", "geometry": "rect"})
        decor_cb("强调条")
        if bullets:
            card_x = "4.7in"
            card_w = "8.1in"
            bullet_paths = _place_bullet_cards(
                sc, out, sid, bullets,
                x=card_x, y_start=0.95, width=card_w,
                accent=accent, tint="#F8FAFC",
                body_color=st.body_color,
                body_size=st.body_size - 2 if has_image else st.body_size,
                card_gap=0.16,
            )
    else:
        # card / minimal 共用头部：顶部色带 + 标题 + 标题竖条
        if "top_bar" in st.decor:
            sc.add(out, sid, {"x": "0in", "y": "0in", "width": "13.33in", "height": "0.15in",
                              "fill": accent, "geometry": "rect"})
            decor_cb("顶部色带")
        officecli.set_element(out, f"{sid}/shape[1]", {
            "x": "1.0in", "y": "0.55in", "width": "11.33in", "height": "0.9in",
            "size": str(st.title_size), "bold": "true", "color": st.title_color, "valign": "top",
        })
        if "title_bar" in st.decor:
            sc.add(out, sid, {"x": "0.55in", "y": "0.7in", "width": "0.13in", "height": "0.85in",
                              "fill": accent, "geometry": "rect"})
            decor_cb("标题竖条")
        if "corner_chip" in st.decor:
            sc.add(out, sid, {"x": "12.55in", "y": "0.55in", "width": "0.45in", "height": "0.45in",
                              "fill": accent, "geometry": "ellipse"})
            decor_cb("角落圆点")

        if bullets:
            # 有背景图时卡片用半透明白底，让背景图透出来；否则用纯白
            card_fill = "#FFFFFFCC" if add_background else "#FFFFFF"
            card_w = "6.6in" if has_image else "11.3in"
            bullet_paths = _place_bullet_cards(
                sc, out, sid, bullets,
                x="1.0in", y_start=1.95, width=card_w,
                accent=accent, tint=tint,
                body_color=st.body_color,
                body_size=st.body_size - 2 if has_image else st.body_size,
                card_gap=0.14,
                card_fill=card_fill,
            )

    # 页脚页码（右下角，主题色）
    p = sc.add(out, sid, {"text": str(slide_num)})
    officecli.set_element(out, p, {
        "x": "12.2in", "y": "7.05in", "width": "0.7in", "height": "0.3in",
        "size": "12", "bold": "true", "color": accent, "align": "right",
    })
    return bullet_paths


def _image_slot(is_cover: bool, st, content_has_image_layout: bool = True) -> dict:
    """配图插入位置（英寸）。按风格布局变体返回不同槽位。"""
    if is_cover:
        if st.cover_layout == "split":
            return {"x": "5.55in", "y": "1.65in", "width": "7.2in", "height": "4.05in"}
        return {"x": "0.9in", "y": "0.55in", "width": "11.53in", "height": "3.6in"}
    if st.content_layout == "sidebar":
        return {"x": "4.7in", "y": "5.05in", "width": "3.9in", "height": "2.19in"}
    if st.content_layout == "minimal":
        return {"x": "7.95in", "y": "1.95in", "width": "4.5in", "height": "2.53in"}
    return {"x": "8.1in", "y": "1.9in", "width": "4.4in", "height": "2.48in"}


async def _call_llm_outline(text: str, want_images: bool) -> tuple[str, list[dict]]:
    """生成大纲并选择风格。返回 (style_name, slides)。"""
    from .ppt_styles import STYLE_CHOICES_BRIEF

    image_rule = ""
    if want_images:
        image_rule = (
            '，每页可有一个可选的 "image_prompt" 字段：仅当该页内容确实需要视觉辅助'
            '（如知识图解、场景示意、对比示意）时给出，纯文字已能表达清楚的页留空；'
            '整份 PPT 配图 2~5 张为宜，不要每页都配。配图描述用中文，60字以内，'
            '描述画面主体、配色与风格，风格统一为"现代扁平信息图插画"，不要包含文字渲染要求'
        )
    system_prompt = (
        "你是专业 PPT 大纲设计师兼视觉总监。用户给出一段文字或一个主题，请完成两件事：\n"
        "A. 从下列预制教育风格中为主题选择最合适的一套：\n" + STYLE_CHOICES_BRIEF + "\n"
        "B. 把内容组织成一份结构清晰、可直接用于演示的 PPT 大纲。\n"
        "要求：\n"
        "1. 每页只聚焦一个主题，页数 4~12 页，内容较多时适当拆分；若用户给的文字本身较长，应充分展开，而不是简略带过。\n"
        "2. 每页包含一个有力的标题（title）、一句话摘要（summary，说明这页要讲什么、方向是什么，30字以内）、"
        "和 4~6 条要点（bullets）" + image_rule + "。\n"
        "3. 用 section 字段把幻灯片分成 2~5 个章节，同一章节的页 section 值相同（如\"背景介绍\"、\"核心方案\"、\"总结展望\"）；"
        "封面页 section 留空。\n"
        "4. 每条要点必须是完整、自成一体的陈述句，能独立说明一个关键信息，尽量具体、有实质内容；\n"
        "   不要用含糊的短语或关键词堆砌，要让观众看着要点就能理解这一页在讲什么。\n"
        "5. 尊重用户原文的信息和语气，不要编造原文没有的事实；可以合理补充过渡性、解释性的说明。\n"
        "6. 中文输出。\n"
        '严格以 JSON 返回：{"style":"风格标识","slides":[{"section":"","title":"...","summary":"...","bullets":["..."],"image_prompt":"可选"}]}。'
        "不要输出任何其他文字。"
    )
    if not settings.browser_use_llm_key:
        raise RuntimeError("未配置 LLM API Key（browser_use_llm_key）")
    url = settings.browser_use_llm_base.rstrip("/") + "/chat/completions"
    payload = {
        "model": settings.browser_use_llm_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": text},
        ],
        "temperature": 0.4,
        "response_format": {"type": "json_object"},
    }
    async with httpx.AsyncClient(timeout=120.0) as client:
        resp = await client.post(
            url,
            headers={"Authorization": f"Bearer {settings.browser_use_llm_key}",
                     "Content-Type": "application/json"},
            json=payload,
        )
        resp.raise_for_status()
        data = resp.json()
    parsed = _safe_json_parse(data["choices"][0]["message"]["content"])
    style_name = ""
    slides_raw: list = []
    if isinstance(parsed, dict):
        style_name = str(parsed.get("style") or "")
        slides_raw = parsed.get("slides") if isinstance(parsed.get("slides"), list) else []
    elif isinstance(parsed, list):
        slides_raw = parsed

    slides = []
    for item in slides_raw:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title", "")).strip()
        bullets = item.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = []
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        if not title and bullets:
            title = bullets.pop(0)
        if not title:
            continue
        image_prompt = str(item.get("image_prompt", "")).strip()
        section = str(item.get("section", "")).strip()
        summary = str(item.get("summary", "")).strip()
        slides.append({
            "section": section,
            "title": title,
            "summary": summary,
            "bullets": bullets,
            "image_prompt": image_prompt,
        })
    if not slides:
        raise RuntimeError("AI 未能生成有效的 PPT 大纲")
    return style_name, slides


async def _vision_polish(out_path: str, slide_num: int, progress_cb) -> None:
    """边写边看：把当前页渲染图发给 vision 模型检查版式，发现明显问题自动微调。

    只修文字溢出/重叠/被截断这类硬伤；vision 未配置或失败时静默跳过。
    """
    if not settings.vision_api_key:
        return
    try:
        img_path = await asyncio.to_thread(take_screenshot, out_path, slide_num, fresh=True)
        with open(img_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
        try:
            os.unlink(img_path)
        except OSError:
            pass

        prompt = (
            "你是 PPT 版式质检员。检查这一页幻灯片截图是否存在明显版式硬伤："
            "文字溢出页面或超出所在色块/卡片、元素互相重叠、文字被截断、字号明显过大或过小。"
            "没有明显问题就返回 {\"ok\": true}。"
            "有问题返回 {\"fixes\": [{\"path\": \"shape[5]\", \"props\": {\"size\": \"14\"}}]}："
            "path 是元素路径（标题为 shape[1]，其余按视觉位置从上到下、从左到右编号 shape[2..N]），"
            "props 只允许 x/y/width/height（英寸字符串如 \"2.0in\"）和 size（纯数字字号）。"
            "最多 3 条修复，只修明显问题，不要 redesign。严格 JSON 输出。"
        )
        url = settings.vision_api_base.rstrip("/") + "/chat/completions"
        async with httpx.AsyncClient(timeout=60.0) as client:
            resp = await client.post(
                url,
                headers={"Authorization": f"Bearer {settings.vision_api_key}",
                         "Content-Type": "application/json"},
                json={
                    "model": settings.vision_model,
                    "messages": [{"role": "user", "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                    ]}],
                    "temperature": 0.1,
                },
            )
            resp.raise_for_status()
            content = resp.json()["choices"][0]["message"]["content"]
        parsed = _safe_json_parse(content)
        if not isinstance(parsed, dict) or parsed.get("ok"):
            return
        fixes = parsed.get("fixes") or []
        applied = 0
        for fix in fixes[:3]:
            if not isinstance(fix, dict):
                continue
            path = str(fix.get("path", ""))
            props = fix.get("props") or {}
            if not re.fullmatch(r"shape\[\d+\]", path) or not isinstance(props, dict):
                continue
            safe_props = {k: str(v) for k, v in props.items() if k in ("x", "y", "width", "height", "size")}
            if not safe_props:
                continue
            try:
                officecli.set_element(out_path, f"/slide[{slide_num}]/{path}", safe_props)
                applied += 1
            except Exception:
                continue
        if applied:
            print(f"[PPT-CREATE-STREAM]   slide[{slide_num}] vision 校验调整 {applied} 处", flush=True)
            progress_cb({"type": "add_decor", "slide": slide_num, "element": "AI 校验调整"})
    except Exception as e:
        print(f"[PPT-CREATE-STREAM]   警告: vision 校验跳过 slide[{slide_num}]: {e}", flush=True)


# ---- 步骤 0b：按文字新建 PPT（流式，逐步放置文字/装饰/图片元素） ----
# 事件：
#   {"type":"outline","slides":[{"title","bullets","image_prompt"}]}
#   {"type":"add_text","slide":N,"element":"title"|"bullet","text":"..."}
#   {"type":"add_decor","slide":N,"element":"..."}              装饰元素放置（色带/卡片/竖条等）
#   {"type":"add_image","slide":N,"status":"generating"|"placed"|"failed"}
#   {"type":"screenshot","slide":N,"image_data":"data:image/png;base64,..."}  关键帧实时渲染
#   {"type":"done", result...} / {"type":"error","message":...}

async def draft_outline(text: str) -> dict:
    """仅生成大纲（含 section/summary/bullets）和推荐风格，不创建 PPT。

    供前端"先审大纲再生成"流程使用：用户可在大纲界面编辑/增删页后，
    再把修改后的 slides 传给 create_from_text_stream。
    """
    want_images = image_gen.is_available()
    style_name, slides = await _call_llm_outline(text, want_images)
    return {"style": style_name, "slides": slides}


async def draft_outline_stream(text: str):
    """流式生成大纲：逐 token 推送 LLM 输出，完成后推送解析后的结构化结果。

    yield 事件：
      {"type": "token", "text": "..."}
      {"type": "done", "style": "...", "slides": [...]}
      {"type": "error", "message": "..."}
    """
    from .ppt_styles import STYLE_CHOICES_BRIEF

    want_images = image_gen.is_available()
    image_rule = ""
    if want_images:
        image_rule = (
            '，每页可有一个可选的 "image_prompt" 字段：仅当该页内容确实需要视觉辅助'
            '（如知识图解、场景示意、对比示意）时给出，纯文字已能表达清楚的页留空；'
            '整份 PPT 配图 2~5 张为宜，不要每页都配。配图描述用中文，60字以内，'
            '描述画面主体、配色与风格，风格统一为"现代扁平信息图插画"，不要包含文字渲染要求'
        )
    system_prompt = (
        "你是专业 PPT 大纲设计师兼视觉总监。用户给出一段文字或一个主题，请完成两件事：\n"
        "A. 从下列预制教育风格中为主题选择最合适的一套：\n" + STYLE_CHOICES_BRIEF + "\n"
        "B. 把内容组织成一份结构清晰、可直接用于演示的 PPT 大纲。\n"
        "要求：\n"
        "1. 每页只聚焦一个主题，页数 4~12 页，内容较多时适当拆分；若用户给的文字本身较长，应充分展开，而不是简略带过。\n"
        "2. 每页包含一个有力的标题（title）、一句话摘要（summary，说明这页要讲什么、方向是什么，30字以内）、"
        "和 4~6 条要点（bullets）" + image_rule + "。\n"
        "3. 用 section 字段把幻灯片分成 2~5 个章节，同一章节的页 section 值相同（如\"背景介绍\"、\"核心方案\"、\"总结展望\"）；"
        "封面页 section 留空。\n"
        "4. 每条要点必须是完整、自成一体的陈述句，能独立说明一个关键信息，尽量具体、有实质内容；\n"
        "   不要用含糊的短语或关键词堆砌，要让观众看着要点就能理解这一页在讲什么。\n"
        "5. 尊重用户原文的信息和语气，不要编造原文没有的事实；可以合理补充过渡性、解释性的说明。\n"
        "6. 中文输出。\n"
        '严格以 JSON 返回：{"style":"风格标识","slides":[{"section":"","title":"...","summary":"...","bullets":["..."],"image_prompt":"可选"}]}。'
        "不要输出任何其他文字。"
    )

    if not settings.browser_use_llm_key:
        yield {"type": "error", "message": "未配置 LLM API Key"}
        return

    url = settings.browser_use_llm_base.rstrip("/") + "/chat/completions"
    payload = {
        "model": settings.browser_use_llm_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": text},
        ],
        "temperature": 0.4,
        "response_format": {"type": "json_object"},
        "stream": True,
    }

    accumulated = ""
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            async with client.stream(
                "POST",
                url,
                headers={
                    "Authorization": f"Bearer {settings.browser_use_llm_key}",
                    "Content-Type": "application/json",
                },
                json=payload,
            ) as resp:
                resp.raise_for_status()
                async for line in resp.aiter_lines():
                    if not line or not line.startswith("data:"):
                        continue
                    data_str = line[5:].strip()
                    if data_str == "[DONE]":
                        break
                    try:
                        chunk = json.loads(data_str)
                        delta = chunk["choices"][0].get("delta", {})
                        token = delta.get("content", "")
                        if token:
                            accumulated += token
                            yield {"type": "token", "text": token}
                    except (json.JSONDecodeError, KeyError, IndexError):
                        continue
    except Exception as e:
        yield {"type": "error", "message": str(e)}
        return

    # 解析完整 JSON
    parsed = _safe_json_parse(accumulated)
    style_name = ""
    slides_raw: list = []
    if isinstance(parsed, dict):
        style_name = str(parsed.get("style") or "")
        slides_raw = parsed.get("slides") if isinstance(parsed.get("slides"), list) else []
    elif isinstance(parsed, list):
        slides_raw = parsed

    slides = []
    for item in slides_raw:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title", "")).strip()
        bullets = item.get("bullets", [])
        if not isinstance(bullets, list):
            bullets = []
        bullets = [str(b).strip() for b in bullets if str(b).strip()]
        if not title and bullets:
            title = bullets.pop(0)
        if not title:
            continue
        slides.append({
            "section": str(item.get("section", "")).strip(),
            "title": title,
            "summary": str(item.get("summary", "")).strip(),
            "bullets": bullets,
            "image_prompt": str(item.get("image_prompt", "")).strip(),
        })

    if not slides:
        yield {"type": "error", "message": "AI 未能生成有效的大纲"}
        return

    yield {"type": "done", "style": style_name, "slides": slides}


def _build_bg_prompt(slide: dict, st) -> str:
    """根据单页大纲生成背景图 prompt：抽象、低对比度、不抢文字注意力。"""
    title = slide.get("title", "")
    summary = slide.get("summary", "")
    section = slide.get("section", "")
    # 从调色板取一个柔和色调描述
    palette_desc = "柔和的蓝紫色调"
    if hasattr(st, "palette") and st.palette:
        accent = st.palette[0][0] if st.palette else "#7C3AED"
        palette_desc = f"以 {accent} 为主色调的柔和渐变"
    return (
        f"PPT背景图，{palette_desc}，抽象几何渐变，低饱和度，"
        f"主题「{title}」{'，' + summary if summary else ''}"
        f"{'，章节：' + section if section else ''}。"
        "画面简洁大气，左侧或边缘有微妙的光影纹理，中央区域保持干净留白以便放置文字，"
        "现代扁平信息图风格，不要出现任何文字，16:9宽屏比例。"
    )


async def create_from_text_stream(text: str, progress_cb: Callable[[dict], None],
                                  style_override: dict | None = None,
                                  slides_override: list[dict] | None = None,
                                  add_background: bool = False) -> dict:
    """根据用户文字生成新 PPT，并实时推送每一步的元素放置过程。

    流程：大纲（AI 自动选风格 / 或用 style_override 指定，如参考 PPT 解析风格）
        → 按需并行生图 → 逐页放置标题/装饰/正文 → vision 边看边改 → 配图，
    每个关键节点立即渲染截图推送，前端可实时看到真实 PPT 画面的变化。
    """
    from .ppt_styles import StyleProfile, get_preset

    print(f"[PPT-CREATE-STREAM] 开始流式新建 PPT: 输入 {len(text)} 字符", flush=True)

    want_images = image_gen.is_available()

    # 1) 大纲：用户已在前端编辑确认则直接使用，否则 AI 生成
    if slides_override:
        slides = []
        for item in slides_override:
            title = str(item.get("title", "")).strip()
            bullets = item.get("bullets", [])
            if not isinstance(bullets, list):
                bullets = []
            bullets = [str(b).strip() for b in bullets if str(b).strip()]
            if not title:
                continue
            slides.append({
                "section": str(item.get("section", "")).strip(),
                "title": title,
                "summary": str(item.get("summary", "")).strip(),
                "bullets": bullets,
                "image_prompt": str(item.get("image_prompt", "")).strip(),
            })
        if not slides:
            raise RuntimeError("大纲为空，请至少保留一页")
        style_name = ""
        print(f"[PPT-CREATE-STREAM] 使用用户确认的大纲: {len(slides)} 页", flush=True)
    else:
        style_name, slides = await _call_llm_outline(text, want_images)

    if style_override:
        st = StyleProfile.from_dict(style_override)
        print(f"[PPT-CREATE-STREAM] 使用参考风格: {st.display_name}", flush=True)
    else:
        st = get_preset(style_name)
        print(f"[PPT-CREATE-STREAM] AI 选择风格: {st.display_name} ({st.name})", flush=True)

    progress_cb({"type": "outline", "slides": slides,
                 "style": {"name": st.name, "display_name": st.display_name}})

    # 2) 大纲就绪后按 image_prompt 并行生图（非每页都有；与文字放置并行）
    img_tasks: dict[int, asyncio.Task] = {}
    bg_tasks: dict[int, asyncio.Task] = {}
    if want_images:
        for i, slide in enumerate(slides):
            ip = slide.get("image_prompt")
            if ip:
                img_tasks[i + 1] = asyncio.create_task(image_gen.generate_image(ip))
            # 背景图：根据每页标题+摘要生成（封面不加，封面已有自己的版式）
            if add_background and i > 0:
                bg_prompt = _build_bg_prompt(slide, st)
                bg_tasks[i + 1] = asyncio.create_task(image_gen.generate_image(bg_prompt, prefix="pptbg"))
        if img_tasks:
            print(f"[PPT-CREATE-STREAM] 并行生图 {len(img_tasks)} 张", flush=True)
        if bg_tasks:
            print(f"[PPT-CREATE-STREAM] 并行生成背景图 {len(bg_tasks)} 张", flush=True)

    work_dir = _get_work_dir()
    out_path = work_dir / f"新建-{int(time.time()*1000)}.pptx"
    print(f"[PPT-CREATE-STREAM] 生成 {len(slides)} 页大纲 -> {out_path}", flush=True)

    if not officecli.create_document(str(out_path)):
        raise RuntimeError("创建新 PPT 失败")

    palette = st.palette
    palette_len = len(palette)

    async def _push_frame(slide_num: int) -> None:
        """关键帧：立即保存并 fresh 截图推送，让用户在编辑阶段就看到真实画面。"""
        try:
            officecli.save_document(str(out_path))
            img_path = await asyncio.to_thread(take_screenshot, str(out_path), slide_num, fresh=True)
            if img_path:
                with open(img_path, "rb") as f:
                    b64 = base64.b64encode(f.read()).decode("ascii")
                progress_cb({
                    "type": "screenshot",
                    "slide": slide_num,
                    "image_data": f"data:image/png;base64,{b64}",
                })
                try:
                    os.unlink(img_path)
                except OSError:
                    pass
        except Exception as e:
            print(f"[PPT-CREATE-STREAM]   警告: 过程截图失败 slide[{slide_num}]: {e}", flush=True)

    # 记录每页配图插入所需的上下文（供第二阶段回填使用）
    slide_ctx: list[dict] = []

    # ── 第一阶段：极速完成所有页面的文字 + 版式（不等图片） ──
    for i, slide in enumerate(slides):
        slide_num = i + 1
        accent, tint = palette[i % palette_len]
        is_cover = slide_num == 1
        has_image = slide_num in img_tasks

        # 添加标题占位符（shape[1]）
        officecli.add_element(str(out_path), "/", "slide", {"title": slide["title"]})
        progress_cb({"type": "add_text", "slide": slide_num, "element": "title", "text": slide["title"]})
        print(
            f"[PPT-CREATE-STREAM]   slide[{slide_num}] 标题: {slide['title']} "
            f"({'封面' if is_cover else '内容'}, accent={accent})",
            flush=True,
        )
        await asyncio.sleep(0.4)  # 让"标题放置"过程可感知

        sid = f"/slide[{slide_num}]"
        bullets = slide["bullets"]

        def _decor(name: str) -> None:
            progress_cb({"type": "add_decor", "slide": slide_num, "element": name})

        body_paths: list[str] = []
        try:
            if is_cover:
                _place_cover(str(out_path), sid, slide, st, accent, tint, has_image, _decor)
            else:
                body_paths = _place_content(str(out_path), sid, slide_num, slide, st,
                                            accent, tint, has_image, _decor, add_background)
        except Exception as e:
            print(f"[PPT-CREATE-STREAM]   警告: 版式设置失败 slide[{slide_num}]: {e}", flush=True)

        await _push_frame(slide_num)  # 关键帧 1：版式骨架就位

        # 逐张填充要点卡片文字（内容页）
        if body_paths and bullets:
            for bp, btext in zip(body_paths, bullets):
                officecli.set_element(str(out_path), bp, {"text": btext})
                progress_cb({"type": "add_text", "slide": slide_num, "element": "bullet", "text": btext})
                await asyncio.sleep(0.35)
            print(f"[PPT-CREATE-STREAM]   slide[{slide_num}] 共 {len(bullets)} 条要点卡片", flush=True)

        await _push_frame(slide_num)  # 关键帧 2：文字全部就位

        # 边写边看：vision 校验本页版式，发现问题自动微调后补一帧
        await _vision_polish(str(out_path), slide_num, progress_cb)

        # 配图仍在后台生成，标记状态后立即继续下一页，不等待
        if has_image:
            progress_cb({"type": "add_image", "slide": slide_num, "status": "generating"})

        # 保存上下文供第二阶段回填图片
        slide_ctx.append({
            "slide_num": slide_num,
            "sid": sid,
            "is_cover": is_cover,
            "accent": accent,
        })

        await asyncio.sleep(0.3)  # 页面之间短暂停顿

    # ── 第二阶段：回填配图 + 背景图（图片已在后台并行生成了一段时间） ──
    # 此时所有文字和版式已完成，用户可以立即看到完整 PPT 骨架；
    # 配图逐张回填，每张插入后刷新该页截图，不影响其他页面的完成度。
    if img_tasks:
        print(f"[PPT-CREATE-STREAM] 文字全部就位，开始回填 {len(img_tasks)} 张配图", flush=True)
    if bg_tasks:
        print(f"[PPT-CREATE-STREAM] 开始回填 {len(bg_tasks)} 张背景图", flush=True)

    for ctx in slide_ctx:
        sn = ctx["slide_num"]

        # ── 背景图（用 picture 元素全幅插入，PowerPoint 中先添加的元素在底层） ──
        if sn in bg_tasks:
            bg_file = None
            try:
                bg_file = await bg_tasks[sn]
            except Exception as e:
                print(f"[PPT-CREATE-STREAM]   警告: 背景图任务异常 slide[{sn}]: {e}", flush=True)
            if bg_file:
                try:
                    # 使用 python-pptx 直接在 XML 层面将图片移到最底层
                    from pptx import Presentation
                    from pptx.util import Inches

                    prs = Presentation(str(out_path))
                    slide = prs.slides[sn - 1]  # 0-indexed
                    left = Inches(0)
                    top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height
                    pic = slide.shapes.add_picture(bg_file, left, top, width, height)
                    # 将图片移到最底层（第一个子元素位置）
                    sp = pic._element
                    sp.getparent().remove(sp)
                    slide.shapes._spTree.insert(2, sp)
                    prs.save(str(out_path))

                    # 添加半透明白色蒙版确保文字可读
                    officecli.add_element(str(out_path), ctx["sid"], "shape", {
                        "x": "0in", "y": "0in", "width": "13.33in", "height": "7.5in",
                        "fill": "#FFFFFF", "opacity": "0.75", "geometry": "rect",
                    })
                    # 将蒙版移到背景图之上、其他内容之下
                    prs2 = Presentation(str(out_path))
                    slide2 = prs2.slides[sn - 1]
                    # 找到刚添加的矩形蒙版（最后一个 shape）
                    shapes_list = list(slide2.shapes)
                    if len(shapes_list) > 1:
                        overlay = shapes_list[-1]._element
                        overlay.getparent().remove(overlay)
                        slide2.shapes._spTree.insert(3, overlay)
                    prs2.save(str(out_path))

                    progress_cb({"type": "add_image", "slide": sn, "status": "placed"})
                    print(f"[PPT-CREATE-STREAM]   slide[{sn}] 背景图已放置", flush=True)
                    await _push_frame(sn)
                except Exception as e:
                    print(f"[PPT-CREATE-STREAM]   警告: 背景图插入失败 slide[{sn}]: {e}", flush=True)

        # ── 内容配图 ──
        if sn in img_tasks:
            img_file = None
            try:
                img_file = await img_tasks[sn]
            except Exception as e:
                print(f"[PPT-CREATE-STREAM]   警告: 生图任务异常 slide[{sn}]: {e}", flush=True)
            if not img_file:
                progress_cb({"type": "add_image", "slide": sn, "status": "failed"})
            else:
                try:
                    slot = _image_slot(ctx["is_cover"], st)
                    officecli.add_element(str(out_path), ctx["sid"], "picture", {"src": img_file, **slot})
                    if not ctx["is_cover"] and st.content_layout == "card":
                        officecli.add_element(str(out_path), ctx["sid"], "shape", {
                            "x": slot["x"], "y": "4.55in", "width": slot["width"], "height": "0.08in",
                            "fill": ctx["accent"], "geometry": "rect",
                        })
                    progress_cb({"type": "add_image", "slide": sn, "status": "placed"})
                    print(f"[PPT-CREATE-STREAM]   slide[{sn}] 配图已放置", flush=True)
                    await _push_frame(sn)
                except Exception as e:
                    print(f"[PPT-CREATE-STREAM]   警告: 配图插入失败 slide[{sn}]: {e}", flush=True)
                    progress_cb({"type": "add_image", "slide": sn, "status": "failed"})

    officecli.save_document(str(out_path))
    final_name = out_path.name
    print(f"[PPT-CREATE-STREAM] 完成: {final_name} ({len(slides)} 页)", flush=True)
    return {
        "file_path": str(out_path),
        "file_name": final_name,
        "total_slides": len(slides),
    }


def build_merged_ppt(sections: list[dict]) -> dict:
    """将各 section 合并为一个总览 PPT。

    返回: {"file_path": str, "file_name": str, "total_slides": int}
    """
    total_parts = sum(len(s["parts"]) for s in sections)
    print(f"[PPT-MERGE] 开始合并: {len(sections)} 个 section, {total_parts} 个文件片段", flush=True)

    work_dir = _get_work_dir()
    temp_path = work_dir / f"merged-{int(time.time()*1000)}.pptx"
    print(f"[PPT-MERGE] 临时合并文件: {temp_path}", flush=True)

    if not officecli.create_document(str(temp_path)):
        raise RuntimeError("创建合并 PPT 失败")

    target_index = 0  # 合并文件当前幻灯片数（0-based 计数，实际 slide[N] 是 1-based）
    failed = 0

    for si, section in enumerate(sections):
        # 插入 section 标题页
        title_ok = officecli.add_element(
            str(temp_path), "/", "slide", {"title": section["name"]}
        )
        if title_ok:
            target_index += 1
            print(f"[PPT-MERGE] 插入 section 标题页 [{si}] \"{section['name']}\" -> target slide {target_index}", flush=True)

        for part in section["parts"]:
            src_path = part["file_path"]
            if not os.path.exists(src_path):
                print(f"[PPT-MERGE] 源文件不存在: {src_path}", flush=True)
                failed += 1
                continue

            for n in range(part["slide_start"], part["slide_end"] + 1):
                src_selector = f"/slide[{n}]"
                try:
                    commands = officecli.dump_element(src_path, src_selector)
                except Exception as e:
                    print(f"[PPT-MERGE] dump 失败: {part['file_name']}{src_selector}: {e}", flush=True)
                    failed += 1
                    continue

                if not commands:
                    failed += 1
                    continue

                target_index += 1
                # 改写路径引用：/slide[src] -> /slide[dst]
                rewritten = _rewrite_slide_paths(commands, n, target_index)

                try:
                    officecli.apply_batch(str(temp_path), rewritten)
                except Exception as e:
                    print(f"[PPT-MERGE] batch 失败: {part['file_name']} slide[{n}]: {e}", flush=True)
                    failed += 1

    # 确保保存
    officecli.save_document(str(temp_path))

    print(f"[PPT-MERGE] 合并完成: {target_index} 页, 失败 {failed} 页", flush=True)

    if target_index == 0:
        try:
            temp_path.unlink(missing_ok=True)
        except Exception:
            pass
        raise RuntimeError("合并结果为空（0 页）")

    final_name = f"合并总览-{int(time.time())}.pptx"
    final_path = work_dir / final_name
    shutil.move(str(temp_path), str(final_path))

    return {
        "file_path": str(final_path),
        "file_name": final_name,
        "total_slides": target_index,
    }


def _rewrite_slide_paths(commands: list[dict], src_index: int, dst_index: int) -> list[dict]:
    """将 dump 命令中的 /slide[src] 引用改写为 /slide[dst]。"""
    src_token = f"/slide[{src_index}]"
    dst_token = f"/slide[{dst_index}]"

    def rewrite(v):
        if isinstance(v, str):
            return v.replace(src_token, dst_token)
        if isinstance(v, list):
            return [rewrite(x) for x in v]
        if isinstance(v, dict):
            return {k: rewrite(val) for k, val in v.items()}
        return v

    return [rewrite(c) for c in commands]


# ---- 步骤 4：AI 统一修改 ----

async def modify_content(sections: list[dict], instruction: str = "") -> list[dict]:
    """AI 统一修改各 PPT 文本内容，返回补丁列表。"""
    patches: list[dict] = []
    all_files = _collect_unique_files(sections)
    print(f"[PPT-MODIFY] 开始统一修改: {len(all_files)} 个文件, 指令: "
          f"{instruction[:80] if instruction else '(默认修正一致性)'}", flush=True)

    overview = build_reading_script(sections)[:12000]

    for fi, file_info in enumerate(all_files):
        file_id = file_info["file_id"]
        file_path = file_info["file_path"]
        file_name = file_info["file_name"]

        # 收集该文件中所有 section 涉及的文本节点
        nodes: list[dict] = []
        for section in sections:
            for part in section["parts"]:
                if part["file_id"] != file_id:
                    continue
                for slide in part["slides"]:
                    for t in slide["texts"]:
                        nodes.append({
                            "section": section["name"],
                            "slide": slide["index"],
                            "path": t["path"],
                            "text": t["text"],
                        })

        if not nodes:
            continue

        print(f"[PPT-MODIFY] [{fi+1}/{len(all_files)}] {file_name} -> {len(nodes)} 个文本节点", flush=True)

        system_prompt = (
            "你是 PPT 内容编辑。用户在合并同一模板的多个 PPT 后，希望对内容做统一修改"
            "（模板一致，无需改版式/字体/配色，只改文字内容）。\n"
            "请对给定文件中的每一条文本节点，在保持原意和结构的前提下，按用户要求修改内容，"
            "并让同一 section 在不同文件间的措辞、数据保持一致。\n"
            "逐条输出，不要漏掉任何节点。严格以 JSON 返回，格式：\n"
            '{"patches":[{"path":"/slide[1]/shape[@id=2]/paragraph[1]","newText":"修改后的文字"}]}\n'
            "path 必须与输入完全一致，newText 为修改后的新文本。未修改的节点 newText 与原文本相同。不要输出其他文字。"
        )

        nodes_text = "\n".join(
            f"[{n['section']}] Slide {n['slide']} | {n['path']} | {n['text']}"
            for n in nodes
        )
        user_content = (
            f"文件：{file_name}\n\n"
            f"用户要求：{instruction or '检查并修正各 section 内容的一致性与明显错误，统一措辞，保持原意。'}\n\n"
            f"跨文件总览参考（用于保持一致）：\n{overview}\n\n"
            f"本文件需修改的文本节点：\n{nodes_text}"
        )

        try:
            raw_patches = await _call_llm_json(system_prompt, user_content, temperature=0.3)
            print(f"[PPT-MODIFY]   AI 返回 {len(raw_patches)} 条补丁", flush=True)

            for rp in raw_patches:
                p_path = rp.get("path", "")
                if not p_path:
                    continue
                node = next((n for n in nodes if n["path"] == p_path), None)
                if not node:
                    continue
                new_text = str(rp.get("newText", "")).strip()
                if new_text == node["text"]:
                    continue
                patches.append({
                    "file_id": file_id,
                    "file_name": file_name,
                    "file_path": file_path,
                    "slide": node["slide"],
                    "path": p_path,
                    "text": node["text"],
                    "new_text": new_text,
                })
        except Exception as e:
            print(f"[PPT-MODIFY] 错误: {file_name} 修改失败: {e}", flush=True)

    print(f"[PPT-MODIFY] 全部完成，累计 {len(patches)} 条补丁", flush=True)
    return patches


def _collect_unique_files(sections: list[dict]) -> list[dict]:
    """从 sections 中收集所有不重复的文件。"""
    seen: dict[str, dict] = {}
    for s in sections:
        for p in s["parts"]:
            if p["file_id"] not in seen:
                seen[p["file_id"]] = {
                    "file_id": p["file_id"],
                    "file_name": p["file_name"],
                    "file_path": p["file_path"],
                }
    return list(seen.values())


# ---- 步骤 5：一键回填 ----

def apply_patches(patches: list[dict]) -> dict:
    """将补丁原位写回各原始 PPT。"""
    applied = 0
    failed = 0
    errors: list[str] = []

    # 按文件分组
    by_file: dict[str, list[dict]] = {}
    for p in patches:
        by_file.setdefault(p["file_path"], []).append(p)

    print(f"[PPT-APPLY] 开始回填: {len(patches)} 条补丁, {len(by_file)} 个文件", flush=True)

    for file_idx, (file_path, file_patches) in enumerate(by_file.items()):
        file_name = file_patches[0]["file_name"]
        if not os.path.exists(file_path):
            failed += len(file_patches)
            errors.append(f"文件不存在: {file_name}")
            continue

        print(f"[PPT-APPLY] [{file_idx+1}/{len(by_file)}] {file_name} -> "
              f"{len(file_patches)} 条补丁", flush=True)

        for i, p in enumerate(file_patches):
            try:
                officecli.set_element(file_path, p["path"], {"text": p["new_text"]})
                applied += 1
            except Exception as e:
                failed += 1
                err_msg = f"{file_name} {p['path']}: {e}"
                errors.append(err_msg)
                print(f"[PPT-APPLY]   失败: {err_msg}", flush=True)

    # 确保所有文件保存
    for file_path in by_file:
        try:
            officecli.save_document(file_path)
        except Exception:
            pass

    print(f"[PPT-APPLY] 回填完成: 成功={applied}, 失败={failed}", flush=True)
    return {"applied": applied, "failed": failed, "errors": errors}


# ---- 截图 ----

# 全局并发信号量：限制最多 3 个 officecli 渲染进程同时运行。
# 已验证：同一文件不同页并发渲染是安全的（每次都是独立进程读取），
# 三个预览窗口首屏可并行出图，总耗时约等于最慢单页而非串行之和。
_screenshot_semaphore = threading.Semaphore(3)


def take_screenshot(file_path: str, page: int = 1, *, fresh: bool = False) -> str:
    """为 PPT 指定页生成截图，返回图片路径。

    优化：
    1. 磁盘缓存：同一文件+页码渲染结果落盘缓存，命中直接返回，避免重复渲染
    2. 全局并发信号量：同文件不同页可并行，首屏多个窗口同时出图
    3. 超时重试，并清理残留临时文件

    fresh=True 时绕过缓存直接渲染且不写缓存——用于生成过程中的实时过程图，
    保证反映当前编辑状态（同页多次截图都能得到最新画面）。
    """
    work_dir = _get_work_dir()
    cache_dir = work_dir / "cache"
    cache_dir.mkdir(parents=True, exist_ok=True)

    # 缓存键：文件路径 + 页码 的哈希
    cache_key = f"{abs(hash(file_path))}-p{page}"
    cache_path = cache_dir / f"{cache_key}.png"

    # 命中缓存：直接返回（fresh 模式跳过）
    if not fresh and cache_path.exists() and cache_path.stat().st_size > 0:
        return str(cache_path)

    # 临时输出路径（渲染时先写临时文件，成功后改名进缓存；fresh 模式直接用临时路径返回）
    tmp_path = cache_dir / f"{cache_key}.{int(time.time()*1000)}.tmp.png"

    last_err: Exception | None = None
    for attempt in range(3):
        try:
            with _screenshot_semaphore:
                officecli.view_screenshot(file_path, str(tmp_path), page=page)
            if tmp_path.exists() and tmp_path.stat().st_size > 0:
                if fresh:
                    # 过程图：不进缓存，直接返回临时路径（调用方读完后可自行清理）
                    return str(tmp_path)
                # 渲染成功，写入缓存
                try:
                    cache_path.unlink(missing_ok=True)
                except OSError:
                    pass
                tmp_path.replace(cache_path)
                return str(cache_path)
            raise RuntimeError("截图文件为空或未生成")
        except subprocess.TimeoutExpired:
            last_err = RuntimeError(f"截图超时（第 {attempt + 1} 次）")
        except Exception as e:
            last_err = e
        # 清理残留临时文件
        if tmp_path.exists():
            try:
                tmp_path.unlink()
            except OSError:
                pass
        if attempt < 2:
            time.sleep(0.5 * (attempt + 1))

    raise RuntimeError(f"截图失败（已重试 3 次）: {last_err}")


# ---- LLM 调用 ----

async def _call_llm_json(system_prompt: str, user_content: str, temperature: float = 0.2) -> list[dict]:
    """调用 LLM 并解析 JSON 数组响应。"""
    if not settings.browser_use_llm_key:
        raise RuntimeError("未配置 LLM API Key（browser_use_llm_key）")

    url = settings.browser_use_llm_base.rstrip("/") + "/chat/completions"
    headers = {
        "Authorization": f"Bearer {settings.browser_use_llm_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": settings.browser_use_llm_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_content},
        ],
        "temperature": temperature,
        "response_format": {"type": "json_object"},
    }

    async with httpx.AsyncClient(timeout=120.0) as client:
        resp = await client.post(url, headers=headers, json=payload)
        resp.raise_for_status()
        data = resp.json()

    content = data["choices"][0]["message"]["content"]
    parsed = _safe_json_parse(content)

    # detect_sections 返回 {"sections": [...]}
    if isinstance(parsed, dict) and "sections" in parsed:
        return parsed["sections"] if isinstance(parsed["sections"], list) else []
    # create_from_text 返回 {"slides": [...]}
    if isinstance(parsed, dict) and "slides" in parsed:
        return parsed["slides"] if isinstance(parsed["slides"], list) else []
    # modify_content 返回 {"patches": [...]}
    if isinstance(parsed, dict) and "patches" in parsed:
        return parsed["patches"] if isinstance(parsed["patches"], list) else []
    # 也可能直接是数组
    if isinstance(parsed, list):
        return parsed
    return []


def _safe_json_parse(text: str) -> Any:
    """安全解析 JSON，容忍 markdown 代码块包裹。"""
    if not text:
        return None
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # 尝试提取 ```json ... ``` 块
    m = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if m:
        try:
            return json.loads(m.group(1))
        except json.JSONDecodeError:
            pass
    # 尝试提取第一个 { ... } 或 [ ... ]
    for start_ch, end_ch in [("{", "}"), ("[", "]")]:
        start = text.find(start_ch)
        end = text.rfind(end_ch)
        if start >= 0 and end > start:
            try:
                return json.loads(text[start:end + 1])
            except json.JSONDecodeError:
                continue
    return None


# ---- 参考 PPT 风格拆解 ----

async def analyze_style(file_path: str,
                        progress_cb: Callable[[dict], None] | None = None) -> dict:
    """拆解参考 PPT 的视觉风格，返回 StyleProfile dict（可直接作 style_override）。

    流程：截取封面 + 最多 2 张内容页 → vision 模型分析配色/字号/版式/装饰
    → 组装 StyleProfile。vision 未配置时抛错（前端应提示先配置识图 AI）。
    """
    from .ppt_styles import StyleProfile

    def _report(stage: str, **extra: Any) -> None:
        if progress_cb:
            try:
                progress_cb({"type": "analyze_style", "stage": stage, **extra})
            except Exception:
                pass

    if not settings.vision_api_key:
        raise RuntimeError("未配置识图 AI（vision_api_key），无法拆解参考风格")

    file_path = str(Path(file_path).resolve())
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    total = get_page_count(file_path)
    # 选代表页：封面 + 中间两个内容页（页少时自动减少）
    candidates = [1]
    for p in (total // 3 or 2, (2 * total) // 3 or 3):
        if 1 < p <= total and p not in candidates:
            candidates.append(p)
    pages = candidates[:3]
    _report("screenshot", pages=pages, total=total)
    print(f"[PPT-STYLE] 拆解参考风格: {Path(file_path).name} 共 {total} 页, 取页 {pages}", flush=True)

    # 并行截图
    images: list[str] = []
    for p in pages:
        img_path = await asyncio.to_thread(take_screenshot, file_path, p)
        with open(img_path, "rb") as f:
            images.append(base64.b64encode(f.read()).decode("ascii"))

    _report("vision", pages=pages)
    prompt = (
        "你是 PPT 视觉风格分析师。给出一份参考 PPT 的若干页截图（第 1 张是封面，其余是内容页），"
        "请拆解它的视觉风格，严格输出如下 JSON：\n"
        "{\n"
        '  "display_name": "风格中文名（4~6字，如 简约蓝白）",\n'
        '  "palette": [["#主色1", "#配套浅色底1"], ["#主色2", "#浅底2"], ...3~4组],\n'
        '  "title_size": 标题字号(整数, 24~40),\n'
        '  "body_size": 正文字号(整数, 12~22),\n'
        '  "title_color": "#标题颜色",\n'
        '  "body_color": "#正文颜色",\n'
        '  "cover_layout": "center(居中式封面) 或 split(左右/上下分块式封面)",\n'
        '  "content_layout": "card(正文在浅色卡片上) / minimal(大留白无卡片) / sidebar(侧边大色块)",\n'
        '  "decor": ["top_bar"(顶部色带), "title_bar"(标题旁竖条), "page_dot"(页码圆点), "corner_chip"(角落色块) 中实际出现的],\n'
        '  "style_notes": "一句话风格描述(中文, 30字以内)"\n'
        "}\n"
        "颜色必须是 7 位十六进制（#RRGGBB）；浅色底与主色同色系、足够浅（接近白）。"
        "只输出 JSON，不要任何其他文字。"
    )
    url = settings.vision_api_base.rstrip("/") + "/chat/completions"
    content_parts: list[dict] = [{"type": "text", "text": prompt}]
    for b64 in images:
        content_parts.append({"type": "image_url",
                              "image_url": {"url": f"data:image/png;base64,{b64}"}})
    async with httpx.AsyncClient(timeout=120.0) as client:
        resp = await client.post(
            url,
            headers={"Authorization": f"Bearer {settings.vision_api_key}",
                     "Content-Type": "application/json"},
            json={
                "model": settings.vision_model,
                "messages": [{"role": "user", "content": content_parts}],
                "temperature": 0.2,
            },
        )
        resp.raise_for_status()
        content = resp.json()["choices"][0]["message"]["content"]

    parsed = _safe_json_parse(content)
    if not isinstance(parsed, dict):
        raise RuntimeError("识图 AI 未返回有效的风格 JSON")
    parsed["name"] = "reference"  # 参考风格统一标识，避免误命中预设
    st = StyleProfile.from_dict(parsed)
    print(f"[PPT-STYLE] 风格拆解完成: {st.display_name} palette={len(st.palette)} 组", flush=True)
    _report("done", style=st.to_dict())
    return st.to_dict()


# ---- 单页元素读取 / 编辑 ----

# 允许前端修改的属性白名单（几何 + 文字 + 外观）
_EDITABLE_PROPS = {
    "x", "y", "width", "height",           # 几何（英寸字符串）
    "text",                                 # 文字内容
    "size", "bold", "italic", "color",     # 字体
    "fill", "align", "valign", "lineSpacing",
}


def _collect_shape_nodes(node: Any, out: list[dict] | None = None) -> list[dict]:
    """递归收集幻灯片中的 shape/picture 节点及其几何与文字属性。"""
    if out is None:
        out = []
    if not node or not isinstance(node, dict):
        return out

    node_type = node.get("type", "")
    if node_type in ("shape", "picture", "graphic-frame", "table"):
        item: dict[str, Any] = {
            "path": node.get("path", ""),
            "type": node_type,
        }
        # 透传常见属性（存在才带）
        for k in ("x", "y", "width", "height", "text", "size", "bold", "italic",
                  "color", "fill", "align", "valign", "geometry", "src", "name"):
            v = node.get(k)
            if v is not None and v != "":
                item[k] = v
        out.append(item)

    for child in node.get("children", []) or []:
        _collect_shape_nodes(child, out)
    return out


def get_slide_elements(file_path: str, slide_num: int) -> dict:
    """获取某页全部元素（供前端拖拽编辑），返回 {"slide": N, "elements": [...]}。"""
    file_path = str(Path(file_path).resolve())
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    node_json = officecli.get_node(file_path, f"/slide[{slide_num}]")
    slide_node = _extract_slide_node(node_json)
    elements = _collect_shape_nodes(slide_node) if slide_node else []
    return {"slide": slide_num, "elements": elements}


def update_slide_elements(file_path: str, slide_num: int, updates: list[dict]) -> dict:
    """批量更新某页元素属性（拖拽位置/双击改文字）。

    updates: [{"path": "shape[2]" 或 "/slide[N]/shape[2]", "props": {...}}]
    只接受白名单属性；更新后保存并失效该页截图缓存。
    """
    file_path = str(Path(file_path).resolve())
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    applied = 0
    errors: list[str] = []
    for u in updates:
        if not isinstance(u, dict):
            continue
        path = str(u.get("path", "")).strip()
        props = u.get("props") or {}
        if not path or not isinstance(props, dict):
            continue
        # 统一为完整路径
        if not path.startswith("/"):
            path = f"/slide[{slide_num}]/{path}"
        safe = {k: str(v) for k, v in props.items() if k in _EDITABLE_PROPS}
        if not safe:
            continue
        try:
            officecli.set_element(file_path, path, safe)
            applied += 1
        except Exception as e:
            errors.append(f"{path}: {e}")

    if applied:
        officecli.save_document(file_path)
        _invalidate_screenshot_cache(file_path, slide_num)
    return {"applied": applied, "errors": errors}


def _invalidate_screenshot_cache(file_path: str, slide_num: int) -> None:
    """删除某页的截图缓存，使下次截图拿到最新画面。"""
    cache_dir = _get_work_dir() / "cache"
    cache_key = f"{abs(hash(file_path))}-p{slide_num}"
    try:
        (cache_dir / f"{cache_key}.png").unlink(missing_ok=True)
    except OSError:
        pass


async def refine_slide(file_path: str, slide_num: int, instruction: str,
                       progress_cb: Callable[[dict], None] | None = None) -> dict:
    """AI 指令只改某一页：读取该页文本 → LLM 生成文本补丁 → 原位写回 → vision 微调。

    返回 {"applied": 补丁数, "image_data": 最新截图 dataURL}。
    """
    def _report(stage: str, **extra: Any) -> None:
        if progress_cb:
            try:
                progress_cb({"type": "refine_slide", "slide": slide_num, "stage": stage, **extra})
            except Exception:
                pass

    file_path = str(Path(file_path).resolve())
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    if not instruction.strip():
        raise ValueError("请输入修改指令")

    # 1) 读取该页文本节点
    _report("read")
    node_json = officecli.get_node(file_path, f"/slide[{slide_num}]")
    slide_node = _extract_slide_node(node_json)
    texts = _collect_text_nodes(slide_node) if slide_node else []
    if not texts:
        raise RuntimeError(f"第 {slide_num} 页没有可修改的文字")

    # 2) LLM 生成文本补丁
    _report("ai")
    system_prompt = (
        "你是 PPT 内容编辑助手。用户给出某一页幻灯片的全部文字节点和一条修改指令，"
        "请只对需要改动的节点输出补丁，保持其余节点不变。\n"
        "要求：尊重原版式结构，不增删节点，只修改文字内容；改写要简洁、适合幻灯片展示；中文输出。\n"
        '严格以 JSON 返回：{"patches":[{"path":"节点path","new_text":"新文字"}]}。'
    )
    user_content = (
        f"第 {slide_num} 页的文字节点：\n"
        + json.dumps(texts, ensure_ascii=False, indent=1)
        + f"\n\n修改指令：{instruction.strip()}"
    )
    patches = await _call_llm_json(system_prompt, user_content, temperature=0.3)

    # 3) 应用补丁
    applied = 0
    for p in patches:
        if not isinstance(p, dict):
            continue
        path = str(p.get("path", "")).strip()
        new_text = p.get("new_text")
        if not path or new_text is None:
            continue
        if not path.startswith("/"):
            path = f"/slide[{slide_num}]/{path}"
        try:
            officecli.set_element(file_path, path, {"text": str(new_text)})
            applied += 1
        except Exception as e:
            print(f"[PPT-REFINE]   警告: 补丁失败 {path}: {e}", flush=True)

    if applied:
        officecli.save_document(file_path)
        _invalidate_screenshot_cache(file_path, slide_num)
        _report("applied", applied=applied)
        # 4) 边写边看：vision 校验改动后的版式
        await _vision_polish(file_path, slide_num, progress_cb or (lambda e: None))
        officecli.save_document(file_path)
        _invalidate_screenshot_cache(file_path, slide_num)

    # 5) 返回最新截图
    img_path = await asyncio.to_thread(take_screenshot, file_path, slide_num, fresh=True)
    with open(img_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("ascii")
    try:
        os.unlink(img_path)
    except OSError:
        pass
    _report("done", applied=applied)
    return {"applied": applied, "image_data": f"data:image/png;base64,{b64}"}
